'''
Author: qiuzx
Date: 2024-12-31 15:21:50
LastEditors: qiuzx
Description: description
'''
import os
from django.http import HttpResponse, FileResponse, JsonResponse
import markdown
from pptx import Presentation
from .base import ApiResponse
from urllib.parse import quote
from rest_framework.views import APIView
from rest_framework.exceptions import APIException

def read_file_content(file_path, file_type):
    """
    读取不同类型的文件并返回文件内容。
    """
    try:
        if file_type in ['doc', 'docx']:
            # 处理 Word 文件
            from docx import Document  # 确保导入 docx 模块
            doc = Document(file_path)
            content = '\n'.join([para.text for para in doc.paragraphs])
            return {'content': content, 'type': 'text'}

        elif file_type == 'xlsx':
            # 处理 Excel 文件
            import pandas as pd  # 确保导入 pandas 模块
            content = pd.read_excel(file_path).to_html()
            return {'content': content, 'type': 'html'}

        elif file_type == 'pptx':
            # 处理 PPT 文件
            presentation = Presentation(file_path)
            slides = []
            for slide in presentation.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        slide_text.append(shape.text)
                slides.append('\n'.join(slide_text))
            content = '\n---\n'.join(slides)
            return {'content': content, 'type': 'text'}

        elif file_type == 'pdf':
            try:
                with open(file_path, 'rb') as f:
                    pdf_content = f.read()  # 读取 PDF 文件的二进制内容
                response = HttpResponse(pdf_content, content_type='application/pdf')  # 设置为 PDF 类型
                response['Content-Disposition'] = f'inline; filename="{os.path.basename(file_path)}"; filename*=UTF-8\'\'{quote(os.path.basename(file_path))}'
                print(response,'dafd')
                return response
            except Exception as e:
                return JsonResponse({'error': str(e)}, status=500)

        elif file_type == 'md':
            # 处理 Markdown 文件
            with open(file_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
            html_content = markdown.markdown(md_content)
            return {'content': html_content, 'type': 'html'}

        else:
            raise APIException(f"Unsupported file type: {file_type}")

    except Exception as e:
        raise APIException(f"Error reading file: {str(e)}")

class FileContent(APIView):
    def post(self, request):
        file_path = request.data.get('file_path')
        base_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '../../../data')
        sanitized_file_path = os.path.normpath(file_path).lstrip(os.sep)
        file_abs_path = os.path.join(base_path, sanitized_file_path)
        if not file_abs_path.startswith(base_path):
            raise ValueError("文件地址不存在")
        if not os.path.exists(file_abs_path):
            raise ValueError("文件不存在")
        file_type = file_path.split('.')[-1].lower()
        print(file_abs_path, file_type, 'xxxxxxxxxx')
        content = read_file_content(file_abs_path, file_type)
        print(content, 'yy')
        
        if isinstance(content, HttpResponse):
            return content
        
        if isinstance(content, FileResponse):
            return content
        
        return ApiResponse.success("", content)
